"""PPTX parser using python-pptx."""

from pptx import Presentation
from app.ingestion.parser import DocumentParser, ParsedDocument, ParsedSection


class PPTXParser(DocumentParser):

    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    async def parse(self, file_path: str, filename: str) -> ParsedDocument:
        prs = Presentation(file_path)
        sections = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Extract tables from slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            texts.append(row_text)

            if texts:
                sections.append(ParsedSection(
                    content="\n".join(texts),
                    page=slide_num,
                    section=f"Slide {slide_num}",
                ))

        return ParsedDocument(
            sections=sections,
            title=filename,
            metadata={"format": "pptx", "slides": len(prs.slides)},
        )
